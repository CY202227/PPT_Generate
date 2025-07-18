"""
Content management tools for PowerPoint MCP Server.
Handles slides, text, images, and content manipulation.
"""
from typing import Dict, List, Optional, Any, Union
from mcp.server.fastmcp import FastMCP
import utils as ppt_utils
import tempfile
import base64
import os
from pptx import Presentation
from openai import OpenAI
import json
import logging
from pptx.util import Pt


def register_content_tools(app: FastMCP, presentations: Dict, get_current_presentation_id, validate_parameters, is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register content management tools with the FastMCP app"""
    
    @app.tool()
    def add_slide(
        layout_index: int = 1,
        title: Optional[str] = None,
        background_type: Optional[str] = None,  # "solid", "gradient", "professional_gradient"
        background_colors: Optional[List[List[int]]] = None,  # For gradient: [[start_rgb], [end_rgb]]
        gradient_direction: str = "horizontal",
        color_scheme: str = "modern_blue",
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Add a new slide to the presentation with optional background styling."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        # Validate layout index
        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            return {
                "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(pres.slide_layouts) - 1}"
            }
        
        try:
            # Add the slide
            slide, layout = ppt_utils.add_slide(pres, layout_index)
            slide_index = len(pres.slides) - 1
            
            # Set title if provided
            if title:
                ppt_utils.set_title(slide, title)
            
            # Apply background if specified
            if background_type == "gradient" and background_colors and len(background_colors) >= 2:
                ppt_utils.set_slide_gradient_background(
                    slide, background_colors[0], background_colors[1], gradient_direction
                )
            elif background_type == "professional_gradient":
                ppt_utils.create_professional_gradient_background(
                    slide, color_scheme, "subtle", gradient_direction
                )
            
            return {
                "message": f"Added slide {slide_index} with layout {layout_index}",
                "slide_index": slide_index,
                "layout_name": layout.name if hasattr(layout, 'name') else f"Layout {layout_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add slide: {str(e)}"
            }

    @app.tool()
    def get_slide_info(slide_index: int, presentation_id: Optional[str] = None) -> Dict:
        """Get information about a specific slide."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            return ppt_utils.get_slide_info(slide, slide_index)
        except Exception as e:
            return {
                "error": f"Failed to get slide info: {str(e)}"
            }

    @app.tool()
    def populate_placeholder(
        slide_index: int,
        placeholder_idx: int,
        text: str,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Populate a placeholder with text."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            ppt_utils.populate_placeholder(slide, placeholder_idx, text)
            return {
                "message": f"Populated placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to populate placeholder: {str(e)}"
            }

    @app.tool()
    def add_bullet_points(
        slide_index: int,
        placeholder_idx: int,
        bullet_points: List[str],
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Add bullet points to a placeholder."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            placeholder = slide.placeholders[placeholder_idx]
            ppt_utils.add_bullet_points(placeholder, bullet_points)
            return {
                "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add bullet points: {str(e)}"
            }

    @app.tool()
    def manage_text(
        slide_index: int,
        operation: str,  # "add", "format", "validate", "format_runs"
        left: float = 1.0,
        top: float = 1.0,
        width: float = 4.0,
        height: float = 2.0,
        text: str = "",
        shape_index: Optional[int] = None,  # For format/validate operations
        text_runs: Optional[List[Dict]] = None,  # For format_runs operation
        # Formatting options
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        underline: Optional[bool] = None,
        color: Optional[List[int]] = None,
        bg_color: Optional[List[int]] = None,
        alignment: Optional[str] = None,
        vertical_alignment: Optional[str] = None,
        # Advanced options
        auto_fit: bool = True,
        validation_only: bool = False,
        min_font_size: int = 8,
        max_font_size: int = 72,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Unified text management tool for adding, formatting, validating text, and formatting multiple text runs."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        # Validate parameters
        validations = {}
        if font_size is not None:
            validations["font_size"] = (font_size, [(is_positive, "must be a positive integer")])
        if color is not None:
            validations["color"] = (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
        if bg_color is not None:
            validations["bg_color"] = (bg_color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
        
        if validations:
            valid, error = validate_parameters(validations)
            if not valid:
                return {"error": error}
        
        try:
            if operation == "add":
                # Add new textbox
                shape = ppt_utils.add_textbox(
                    slide, left, top, width, height, text,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    color=tuple(color) if color else None,
                    bg_color=tuple(bg_color) if bg_color else None,
                    alignment=alignment,
                    vertical_alignment=vertical_alignment,
                    auto_fit=auto_fit
                )
                return {
                    "message": f"Added text box to slide {slide_index}",
                    "shape_index": len(slide.shapes) - 1,
                    "text": text
                }
            
            elif operation == "format":
                # Format existing text shape
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for formatting: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                shape = slide.shapes[shape_index]
                ppt_utils.format_text_advanced(
                    shape,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    color=tuple(color) if color else None,
                    bg_color=tuple(bg_color) if bg_color else None,
                    alignment=alignment,
                    vertical_alignment=vertical_alignment
                )
                return {
                    "message": f"Formatted text shape {shape_index} on slide {slide_index}"
                }
            
            elif operation == "validate":
                # Validate text fit
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for validation: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                validation_result = ppt_utils.validate_text_fit(
                    slide.shapes[shape_index],
                    text_content=text or None,
                    font_size=font_size or 12
                )
                
                if not validation_only and validation_result.get("needs_optimization"):
                    # Apply automatic fixes
                    fix_result = ppt_utils.validate_and_fix_slide(
                        slide,
                        auto_fix=True,
                        min_font_size=min_font_size,
                        max_font_size=max_font_size
                    )
                    validation_result.update(fix_result)
                
                return validation_result
            
            elif operation == "format_runs":
                # Format multiple text runs with different formatting
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for format_runs: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                if not text_runs:
                    return {"error": "text_runs parameter is required for format_runs operation"}
                
                shape = slide.shapes[shape_index]
                
                # Check if shape has text
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    return {"error": "Shape does not contain text"}
                
                # Clear existing text and rebuild with formatted runs
                text_frame = shape.text_frame
                text_frame.clear()
                
                formatted_runs = []
                
                for run_data in text_runs:
                    if 'text' not in run_data:
                        continue
                        
                    # Add paragraph if needed
                    if not text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()
                    
                    # Add run with text
                    run = paragraph.add_run()
                    run.text = run_data['text']
                    
                    # Apply formatting using pptx imports
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor
                    
                    if 'bold' in run_data:
                        run.font.bold = run_data['bold']
                    if 'italic' in run_data:
                        run.font.italic = run_data['italic']
                    if 'underline' in run_data:
                        run.font.underline = run_data['underline']
                    if 'font_size' in run_data:
                        run.font.size = Pt(run_data['font_size'])
                    if 'font_name' in run_data:
                        run.font.name = run_data['font_name']
                    if 'color' in run_data and is_valid_rgb(run_data['color']):
                        run.font.color.rgb = RGBColor(*run_data['color'])
                    if 'hyperlink' in run_data:
                        run.hyperlink.address = run_data['hyperlink']
                    
                    formatted_runs.append({
                        "text": run_data['text'],
                        "formatting_applied": {k: v for k, v in run_data.items() if k != 'text'}
                    })
                
                return {
                    "message": f"Applied formatting to {len(formatted_runs)} text runs on shape {shape_index}",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "formatted_runs": formatted_runs
                }
            
            else:
                return {
                    "error": f"Invalid operation: {operation}. Must be 'add', 'format', 'validate', or 'format_runs'"
                }
        
        except Exception as e:
            return {
                "error": f"Failed to {operation} text: {str(e)}"
            }

    @app.tool()
    def manage_image(
        slide_index: int,
        operation: str,  # "add", "enhance"
        image_source: str,  # file path or base64 string
        source_type: str = "file",  # "file" or "base64"
        left: float = 1.0,
        top: float = 1.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
        # Enhancement options
        enhancement_style: Optional[str] = None,  # "presentation", "custom"
        brightness: float = 1.0,
        contrast: float = 1.0,
        saturation: float = 1.0,
        sharpness: float = 1.0,
        blur_radius: float = 0,
        filter_type: Optional[str] = None,
        output_path: Optional[str] = None,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Unified image management tool for adding and enhancing images."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            if operation == "add":
                if source_type == "base64":
                    # Handle base64 image
                    try:
                        image_data = base64.b64decode(image_source)
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                            temp_file.write(image_data)
                            temp_path = temp_file.name
                        
                        # Add image from temporary file
                        shape = ppt_utils.add_image(slide, temp_path, left, top, width, height)
                        
                        # Clean up temporary file
                        os.unlink(temp_path)
                        
                        return {
                            "message": f"Added image from base64 to slide {slide_index}",
                            "shape_index": len(slide.shapes) - 1
                        }
                    except Exception as e:
                        return {
                            "error": f"Failed to process base64 image: {str(e)}"
                        }
                else:
                    # Handle file path
                    if not os.path.exists(image_source):
                        return {
                            "error": f"Image file not found: {image_source}"
                        }
                    
                    shape = ppt_utils.add_image(slide, image_source, left, top, width, height)
                    return {
                        "message": f"Added image to slide {slide_index}",
                        "shape_index": len(slide.shapes) - 1,
                        "image_path": image_source
                    }
            
            elif operation == "enhance":
                # Enhance existing image file
                if source_type == "base64":
                    return {
                        "error": "Enhancement operation requires file path, not base64 data"
                    }
                
                if not os.path.exists(image_source):
                    return {
                        "error": f"Image file not found: {image_source}"
                    }
                
                if enhancement_style == "presentation":
                    # Apply professional enhancement
                    enhanced_path = ppt_utils.apply_professional_image_enhancement(
                        image_source, style="presentation", output_path=output_path
                    )
                else:
                    # Apply custom enhancement
                    enhanced_path = ppt_utils.enhance_image_with_pillow(
                        image_source,
                        brightness=brightness,
                        contrast=contrast,
                        saturation=saturation,
                        sharpness=sharpness,
                        blur_radius=blur_radius,
                        filter_type=filter_type,
                        output_path=output_path
                    )
                
                return {
                    "message": f"Enhanced image: {image_source}",
                    "enhanced_path": enhanced_path
                }
            
            else:
                return {
                    "error": f"Invalid operation: {operation}. Must be 'add' or 'enhance'"
                }
        
        except Exception as e:
            return {
                "error": f"Failed to {operation} image: {str(e)}"
            }

def generate_content_from_outline(
    template_path: str,
    outline_path: str,
    output_path: str,
    api_key: str,
    base_url: str,
    model_name: str
) -> Dict[str, Any]:
    """
    根据模板和大纲生成PPT内容
    
    Args:
        template_path: PPT模板文件路径
        outline_path: 大纲JSON文件路径
        output_path: 输出PPT文件路径
        api_key: OpenAI API密钥
        base_url: OpenAI API基础URL（可选）
        model_name: 使用的模型名称（可选，默认为Qwen-72B）
    
    Returns:
        Dict[str, Any]: 包含生成结果的字典
    """
    try:
        # 初始化OpenAI客户端
        client = OpenAI(api_key=api_key, base_url=base_url)
        
        # 加载大纲文件
        with open(outline_path, 'r', encoding='utf-8') as f:
            outline = json.load(f)
            
        # 加载PPT模板
        prs = Presentation(template_path)
        
        # 遍历大纲中的每个部分
        for slide_index, slide in enumerate(prs.slides):
            # 获取当前部分的大纲内容
            section_index = min(slide_index // 5, len(outline["sections"]) - 1)
            current_section = outline["sections"][section_index]
            
            # 为每个形状生成内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 分析形状的样式和位置
                    shape_info = {
                        "type": "text",
                        "position": {"x": shape.left, "y": shape.top},
                        "size": {"width": shape.width, "height": shape.height}
                    }
                    
                    # 使用OpenAI生成内容
                    prompt = f"""
                    根据以下大纲信息生成PPT内容：
                    {json.dumps(current_section, ensure_ascii=False)}
                    
                    形状信息：
                    {json.dumps(shape_info, ensure_ascii=False)}
                    
                    请生成适合该形状的专业PPT内容。
                    """
                    
                    response = client.chat.completions.create(
                        model=model_name,
                        messages=[
                            {"role": "system", "content": "你是一个专业的PPT内容生成助手。"},
                            {"role": "user", "content": prompt}
                        ],
                        max_tokens=150
                    )
                    
                    # 更新形状的文本内容
                    generated_text = response.choices[0].message.content.strip()
                    shape.text_frame.text = generated_text
                    
                    # 应用专业样式
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(18)  # 设置默认字号
                        paragraph.font.name = 'Microsoft YaHei'  # 设置默认字体
        
        # 保存生成的PPT
        prs.save(output_path)
        
        return {
            "success": True,
            "message": "PPT内容生成成功",
            "output_path": output_path
        }
        
    except Exception as e:
        logging.error(f"生成PPT内容失败: {e}")
        return {
            "success": False,
            "error": str(e)
        }